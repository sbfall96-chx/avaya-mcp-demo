import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch

# =========================================================================
# PART 1: POWERPOINT GENERATOR (python-pptx)
# =========================================================================
def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors (Dark Theme matching the UI)
    DARK_NAVY = RGBColor(10, 14, 23)
    LIGHT_GRAY = RGBColor(243, 244, 246)
    TEXT_MUTED = RGBColor(156, 163, 175)
    AVAYA_RED = RGBColor(218, 41, 28)
    CYAN_GLOW = RGBColor(0, 240, 255)
    GREEN_ROI = RGBColor(16, 185, 129)
    YELLOW_WARN = RGBColor(245, 158, 11)

    def apply_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_slide_header(slide, title, subtitle=None):
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.8), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Outfit"
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(14)
            p2.font.color.rgb = CYAN_GLOW
            p2.font.name = "Outfit"
            p2.space_before = Pt(4)

    # Slide 1: Title (Interview Intro)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, DARK_NAVY)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "AVAYA INFINITY MCP CONSOLE"
    p1.font.bold = True
    p1.font.size = Pt(46)
    p1.font.color.rgb = AVAYA_RED
    p1.font.name = "Outfit"

    p2 = tf.add_paragraph()
    p2.text = "Solutions Consultant Technical Demo & Presentation"
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Outfit"
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "A technical deep-dive into why I built this application, what it shows under the hood,\nand how it directly demonstrates the business value of Avaya's actual Infinity roadmap."
    p3.font.size = Pt(15)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Outfit"
    p3.space_before = Pt(30)

    p4 = tf.add_paragraph()
    p4.text = "Presented by: Solutions Consultant Candidate"
    p4.font.bold = True
    p4.font.size = Pt(14)
    p4.font.color.rgb = CYAN_GLOW
    p4.font.name = "Outfit"
    p4.space_before = Pt(40)

    # Slide 2: Why I Built This Application (The Intent)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide2, DARK_NAVY)
    add_slide_header(slide2, "Why I Built This Application: SC Perspective", "Bridging Abstract Architecture with Quantifiable C-Suite Value")
    
    # Left Column
    leftBox = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l = leftBox.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "The Engineering & Demo Imperative"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Outfit"
    
    bullets = [
      "To 'Show, Not Just Tell': Standardized protocol frameworks like Model Context Protocol (MCP) and Lakehouses are abstract concepts. I built this console to make them visual, interactive, and immediately understandable to C-level buyers.",
      "Sales Enablement Focused: Created a client-ready playground that functions offline (perfect for situations with spotty network coverage) or connects directly to live AI APIs for real tool-calling validation.",
      "Realistic Scenario Sandbox: Includes pre-configured incident presets (wait-time spikes, churn escalations) that sales engineering teams can run in real-time to walk clients through day-in-the-life contact center challenges."
    ]
    for b in bullets:
        bp = tf_l.add_paragraph()
        bp.text = "• " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Right Column
    rightBox = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_r = rightBox.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "Interview Alignment & Technical Rigor"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = CYAN_GLOW
    p.font.name = "Outfit"
    
    bullets_r = [
      "Demonstrating SC Competencies: Proves my ability to architect a full-stack dashboard, model complex database nodes, implement visual analytics (Chart.js), and code responsive CSS layouts.",
      "Strategic Solution Strategy: Directly showcases my grasp of Avaya's product portfolio—highlighting the value proposition of AXP cloud telemetry and ACM core media layers.",
      "Interactive Value Pitch: Designed specifically to address Carlos and other sales leaders' target needs—focusing on ROI calculators, compliance controls, and SLA penalty mitigations."
    ]
    for b in bullets_r:
        bp = tf_r.add_paragraph()
        bp.text = "✔ " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Slide 3: What the Application Actually Shows
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide3, DARK_NAVY)
    add_slide_header(slide3, "What the Application Actually Shows", "An Interactive Guided Tour of the Core Modules Built")
    
    # Grid Columns (4 Pillars)
    col_width = Inches(5.5)
    row_height = Inches(2.2)
    
    # Top Left
    tl = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), col_width, row_height)
    tf_tl = tl.text_frame
    tf_tl.word_wrap = True
    p = tf_tl.paragraphs[0]
    p.text = "1. Live Protocol Trace (The 'How')"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = CYAN_GLOW
    p.font.name = "Outfit"
    p_desc = tf_tl.add_paragraph()
    p_desc.text = "Provides step-by-step transparency. It captures raw JSON inputs, tool schemas, and output payloads exchanged between the model client and MCP database servers, demystifying AI decision logic."
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.font.name = "Outfit"
    p_desc.space_before = Pt(4)

    # Top Right
    tr = slide3.shapes.add_textbox(Inches(6.8), Inches(1.8), col_width, row_height)
    tf_tr = tr.text_frame
    tf_tr.word_wrap = True
    p = tf_tr.paragraphs[0]
    p.text = "2. Unified Lakehouse & Silos (The 'What')"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = YELLOW_WARN
    p.font.name = "Outfit"
    p_desc = tf_tr.add_paragraph()
    p_desc.text = "Visualizes simulated CRM profiles, SQL database logs, and SharePoint Knowledge Base tables. Proves that models can retrieve siloed customer data without copying or duplicating data repositories."
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.font.name = "Outfit"
    p_desc.space_before = Pt(4)

    # Bottom Left
    bl = slide3.shapes.add_textbox(Inches(0.75), Inches(4.3), col_width, row_height)
    tf_bl = bl.text_frame
    tf_bl.word_wrap = True
    p = tf_bl.paragraphs[0]
    p.text = "3. Interactive Network Playground (The 'Why')"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN_ROI
    p.font.name = "Outfit"
    p_desc = tf_bl.add_paragraph()
    p_desc.text = "Animates active connection lines and flowing data packets. Instantly models legacy integration mess (N x M spaghetti) vs the clean unified bus (N + M MCP) and Edify policy routing layouts."
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.font.name = "Outfit"
    p_desc.space_before = Pt(4)

    # Bottom Right
    br = slide3.shapes.add_textbox(Inches(6.8), Inches(4.3), col_width, row_height)
    tf_br = br.text_frame
    tf_br.word_wrap = True
    p = tf_br.paragraphs[0]
    p.text = "4. Databricks Governance Log (The 'Security')"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = AVAYA_RED
    p.font.name = "Outfit"
    p_desc = tf_br.add_paragraph()
    p_desc.text = "Logs compliance activities in real-time. Highlights PII redactions, RBAC access blocks, and logs all gateway actions directly into Unity Catalog audit tables for total transparency."
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.font.name = "Outfit"
    p_desc.space_before = Pt(4)

    # Slide 4: Aligning with Avaya's Actual Roadmap
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide4, DARK_NAVY)
    add_slide_header(slide4, "How It Represents Avaya's Actual Infinity", "Translating Architecture Guidelines to Interactive Solutions")
    
    # Left Column
    leftBox4 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l4 = leftBox4.text_frame
    tf_l4.word_wrap = True
    
    p = tf_l4.paragraphs[0]
    p.text = "Core Infrastructure Pillars"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Outfit"
    
    bullets = [
      "AI-Agnostic BYOAI Framework: In the UI, the operator can switch AI models seamlessly. This models Avaya's commitment to avoiding model lock-in, letting customers bring any LLM client to their data.",
      "Zero-Copy Data Access: Under the hood, the console queries database tables dynamically via schema tool definitions. Data remains in-place, aligned with Avaya's zero-copy architecture roadmap.",
      "Edify Journey Flowchart: An animated flowchart diagram mapping dynamic CX stream packet routes, illustrating how Avaya Infinity handles no-code policy orchestration."
    ]
    for b in bullets:
        bp = tf_l4.add_paragraph()
        bp.text = "• " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Right Column
    rightBox4 = slide4.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_r4 = rightBox4.text_frame
    tf_r4.word_wrap = True
    
    p = tf_r4.paragraphs[0]
    p.text = "Operational Controls & Telemetry"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = CYAN_GLOW
    p.font.name = "Outfit"
    
    bullets_r = [
      "Tandem Care Checkpoint: Triggers a visual checkpoint requiring supervisor authorization for high-stakes actions, reflecting Avaya's human-in-the-loop safety roadmap.",
      "Multi-Layer Telemetry Coexistence: Visualizes cloud events (AXP) and voice trunk parameters (ACM Jitter/Loss/MOS) side-by-side. Demonstrates Avaya's capability to bridge legacy telephony cores with cloud operations.",
      "Databricks Unity Catalog Alignment: Incorporates PII HIPAA masking and RBAC query block triggers, proving enterprise-grade compliance safeguards."
    ]
    for b in bullets_r:
        bp = tf_r4.add_paragraph()
        bp.text = "✔ " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Slide 5: Deep Dive: Unified Lakehouse ( SC walkthrough)
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide5, DARK_NAVY)
    add_slide_header(slide5, "Detailed Walkthrough: Unified Lakehouse", "Silo Consolidation & Interactive Data Management Features")
    
    # Left Column
    leftBox5 = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l5 = leftBox5.text_frame
    tf_l5.word_wrap = True
    
    p = tf_l5.paragraphs[0]
    p.text = "Why Silos Are in the Console"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Outfit"
    
    bullets = [
      "Simulating Real-World Friction: Enterprises have customer history, ticket CRM data, and policies separated across different database platforms. This view shows them all in one place.",
      "Zero-Copy Verification: Rather than migrating databases to the LLM, the model issues on-the-fly read queries to separate tables, illustrating zero-copy architectural advantages.",
      "Interactive SharePoint Document Viewer: Users can select actual warranty, tax, or outage documents from a visual index table and immediately read the text clauses used by the AI model."
    ]
    for b in bullets:
        bp = tf_l5.add_paragraph()
        bp.text = "• " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Right Column
    rightBox5 = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_r5 = rightBox5.text_frame
    tf_r5.word_wrap = True
    
    p = tf_r5.paragraphs[0]
    p.text = "Dynamic Data Silos Features"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = CYAN_GLOW
    p.font.name = "Outfit"
    
    bullets_r = [
      "Unified Tables Index: Includes PostgreSQL Call Database, Salesforce/Zendesk CRM Records, and SharePoint Document Libraries.",
      "Interactive Past Interactions Logger: Let's the presenter manually write and submit a supervisor log entry back to the customer's CRM history, showcasing active writebacks.",
      "Dynamic Customer Seeding: Adding a new CRM customer profile automatically seeds historical interaction records based on account tiers (VIP vs Basic)."
    ]
    for b in bullets_r:
        bp = tf_r5.add_paragraph()
        bp.text = "✔ " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Slide 6: Deep Dive: Security & Governance
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide6, DARK_NAVY)
    add_slide_header(slide6, "Detailed Walkthrough: Security & Governance", "Centralized Policy Enforcement, HIPAA Protections, and Real-Time Auditing")
    
    # Left Column
    leftBox6 = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l6 = leftBox6.text_frame
    tf_l6.word_wrap = True
    
    p = tf_l6.paragraphs[0]
    p.text = "Governance Controls & Enforcements"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Outfit"
    
    bullets = [
      "Dynamic PII Masking: Automatically redacts sensitive names (e.g. Jeff Edwards -> J*** E******) in raw database tables, JSON parameters, trace payloads, and final chat outputs.",
      "Centralized Proxy Redaction: Demonstrates that data is protected at the API proxy layer, eliminating the risk of LLM prompts leaking sensitive information.",
      "Restricted DB Access Gatekeeping: Toggling the restriction blocks SQL execution instantly, demonstrating Unity Catalog role-based checks yielding 403 Access Denied messages in the trace."
    ]
    for b in bullets:
        bp = tf_l6.add_paragraph()
        bp.text = "• " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Right Column
    rightBox6 = slide6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_r6 = rightBox6.text_frame
    tf_r6.word_wrap = True
    
    p = tf_r6.paragraphs[0]
    p.text = "Governance Audit Dashboard Features"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = CYAN_GLOW
    p.font.name = "Outfit"
    
    bullets_r = [
      "PII Masked Counter: Running tracker tallying every redaction event occurred in active simulation pipelines.",
      "Blocked Requests Counter: Real-time count of blocked unauthorized calls, highlighting active security gatekeeping.",
      "Compliance SLA Status Indicator: Automatically shows 100% SLA Health, which shifts dynamically if unauthorized query violations exceed safe operating boundaries.",
      "Auditable Logs: Standardized log format matching Databricks catalog traces, making compliance checks simple."
    ]
    for b in bullets_r:
        bp = tf_r6.add_paragraph()
        bp.text = "✔ " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Slide 7: Business Outcomes & Value Case
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide7, DARK_NAVY)
    add_slide_header(slide7, "Proving the Business Outcome & ROI Case", "Translating Telemetry Metrics into Financial Outcomes for Sales Leadership")
    
    # Left Column
    leftBox7 = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l7 = leftBox7.text_frame
    tf_l7.word_wrap = True
    
    p = tf_l7.paragraphs[0]
    p.text = "The Architecture Playground ROI Calculator"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Outfit"
    
    bullets = [
      "Quantified Cost Reduction: Dynamic calculator card showing Legacy custom connectors ($150,000/yr TCO) vs Avaya Infinity MCP Bus ($15,000/yr TCO)—a 90% maintenance cost savings.",
      "Agility Metric: Compares 4 months deployment cycles for point-to-point custom adapters against 2 hours for standardized MCP client-server registrations.",
      "Security Policy Scope: Measures the compliance exposure difference between Legacy (High Risk endpoint leaks) and Avaya Infinity (Protected unified gateway)."
    ]
    for b in bullets:
        bp = tf_l7.add_paragraph()
        bp.text = "• " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Right Column
    rightBox7 = slide7.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_r7 = rightBox7.text_frame
    tf_r7.word_wrap = True
    
    p = tf_r7.paragraphs[0]
    p.text = "The Aura AI Business Outcome Banners"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = CYAN_GLOW
    p.font.name = "Outfit"
    
    bullets_r = [
      "Queue Performance (Preset 1): Dynamic breakout rescheduling and automated chat diversions prevent $3,200/day in SLA fines.",
      "Contract Security (Preset 2): Priority routing VIP callers during jitter spikes retains $24,000 in VIP account contract ARR.",
      "Revenue Retention (Preset 5): Leverages policy guidelines to authorize a $1,260 coupon discount to save an $8,400 churn contract, achieving 6.6x ROI on CAC."
    ]
    for b in bullets_r:
        bp = tf_r7.add_paragraph()
        bp.text = "✔ " + b
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_MUTED
        bp.font.name = "Outfit"
        bp.space_before = Pt(10)

    # Slide 8: SC Pitch Closing
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide8, DARK_NAVY)
    add_slide_header(slide8, "Summary: A Complete Sales Pitch Delivery", "Aligning Customer Success, Security Compliance, and IT Budgets")
    
    txBox = slide8.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(4.5))
    tf8 = txBox.text_frame
    tf8.word_wrap = True

    p1 = tf8.paragraphs[0]
    p1.text = "Why I Am Prepared to Stand as an Avaya Solutions Consultant:"
    p1.font.bold = True
    p1.font.size = Pt(20)
    p1.font.color.rgb = CYAN_GLOW
    p1.font.name = "Outfit"

    points = [
      "1. Ability to Simplify Complexity: Translated dry protocol specifications (MCP schemas, JSON payloads) into a high-fidelity visual experience that any buyer can easily understand.",
      "2. Outcomes-Driven Demos: Put financial results and TCO metrics front and center, ensuring that the technology is directly coupled with the client's business goals ($ savings, SLA compliance, ARR retention).",
      "3. Mastery of the Portfolio: Handcrafted an application representing Avaya's core architectural tenets—proving I have the technical depth to consult clients on actual roadmaps."
    ]
    for pt in points:
        bp = tf8.add_paragraph()
        bp.text = pt
        bp.font.size = Pt(14.5)
        bp.font.color.rgb = LIGHT_GRAY
        bp.font.name = "Outfit"
        bp.space_before = Pt(14)

    # Save
    prs.save("Avaya_Infinity_MCP_Pitch_Deck.pptx")
    print("PowerPoint deck generated: Avaya_Infinity_MCP_Pitch_Deck.pptx")


# =========================================================================
# PART 2: PDF BRIEF GENERATOR (ReportLab)
# =========================================================================
def generate_pdf():
    pdf_filename = "Avaya_Infinity_MCP_Executive_Brief.pdf"
    doc = SimpleDocTemplate(
        pdf_filename,
        pagesize=letter,
        rightMargin=54,
        leftMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=24,
        leading=28,
        textColor=colors.HexColor('#DA291C'),
        spaceAfter=6
    )
    
    subtitle_style = ParagraphStyle(
        'DocSubTitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=12,
        leading=16,
        textColor=colors.HexColor('#6b7280'),
        spaceAfter=15
    )
    
    h2_style = ParagraphStyle(
        'DocH2',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=14,
        leading=18,
        textColor=colors.HexColor('#111827'),
        spaceBefore=12,
        spaceAfter=6,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#374151'),
        spaceAfter=8
    )
    
    bullet_style = ParagraphStyle(
        'DocBullet',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#374151'),
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=4
    )
    
    bold_body_style = ParagraphStyle(
        'DocBoldBody',
        parent=body_style,
        fontName='Helvetica-Bold'
    )
    
    story = []
    
    # ----------------- SECTION 1 -----------------
    story.append(Paragraph("AVAYA INFINITY SOLUTIONS BRIEF: SC INTERVIEW GUIDE", title_style))
    story.append(Paragraph("A Technical Deep-Dive into the Architecture, Security, and Business Value of the MCP Prototype", subtitle_style))
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("1. Purpose & Application Context", h2_style))
    story.append(Paragraph(
        "This Solutions Brief serves as a technical walkthrough designed to present during solutions consulting and "
        "sales engineering interview cycles. The accompanying console application was created to bridge the gap between "
        "abstract database standards—such as the Model Context Protocol (MCP)—and the tangible, high-level business case "
        "metrics that C-suite stakeholders prioritize ($ TCO savings, deployment speed, and HIPAA SLA protection).",
        body_style
    ))
    story.append(Paragraph(
        "By presenting this application, the candidate demonstrates an understanding of Avaya's actual Infinity roadmap: "
        "including AI-agnostic models integration (BYOAI), zero-copy Lakehouses, human-in-the-loop priority routing safeguards "
        "(Tandem Care), and strict zero-trust audit compliance checked by Databricks Unity Catalog.",
        body_style
    ))
    
    story.append(Spacer(1, 10))
    story.append(Paragraph("2. Financial Value: Legacy Spaghetti vs. Avaya Infinity", h2_style))
    story.append(Paragraph(
        "The Architecture Playground tab in the application models the direct business impact of simplifying integration pipelines "
        "from legacy custom-coded connectors ($N \\times M$ paths) to the standardized Model Context Protocol ($N + M$ bus):",
        body_style
    ))
    
    # TCO Table
    table_data = [
        [Paragraph("Metric", bold_body_style), 
         Paragraph("Legacy Point-to-Point Connectors", bold_body_style), 
         Paragraph("Avaya Infinity MCP Bus Architecture", bold_body_style)],
        [Paragraph("Annual TCO", body_style), Paragraph("$150,000 / year in developer maintenance", body_style), Paragraph("$15,000 / year (90% maintenance cost savings)", body_style)],
        [Paragraph("Deployment Speed", body_style), Paragraph("4 Months per enterprise data repository", body_style), Paragraph("2 Hours plug-and-play client-server setup", body_style)],
        [Paragraph("Security Policy Scope", body_style), Paragraph("High Risk (Static files, scattered API tokens)", body_style), Paragraph("Protected (Centralized zero-trust API gateway)", body_style)],
        [Paragraph("Integration Bus", body_style), Paragraph("N x M bespoke integration wrappers", body_style), Paragraph("N + M unified schemas exposed to models", body_style)]
    ]
    
    t = Table(table_data, colWidths=[1.5*inch, 2.5*inch, 2.5*inch])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f3f4f6')),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('BOTTOMPADDING', (0,0), (-1,0), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#d1d5db')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 10))
    
    # ----------------- SECTION 3 -----------------
    story.append(Paragraph("3. Databricks Unity Catalog Enterprise Governance", h2_style))
    story.append(Paragraph(
        "Avaya Infinity enforces zero-trust security parameters directly through the gateway database layer:",
        body_style
    ))
    story.append(Paragraph("• <b>Dynamic HIPAA Masking:</b> The proxy layer intercepts customer records (like Jeff Edwards) and replaces them with redacted formats (J*** E******) in payloads, SQL INSERT events, and chat windows.", bullet_style))
    story.append(Paragraph("• <b>Centralized RBAC Rules:</b> Unauthorized database queries are blocked at the gateway level, returning 403 Security Exceptions, preventing data leakage.", bullet_style))
    story.append(Paragraph("• <b>Active Compliance Dashboard:</b> Running counters monitor PII items masked, request blocks, and SLA compliance health in real-time.", bullet_style))
    
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("4. Presets & Simulated Business Outcomes", h2_style))
    
    scenarios_data = [
        [Paragraph("Scenario Preset", bold_body_style), Paragraph("Aura AI Tool Action", bold_body_style), Paragraph("Quantified Outcome & ROI", bold_body_style)],
        [
            Paragraph("<b>Queue Recovery Spike</b>", body_style), 
            Paragraph("Queries active schedules via Postgres & Zendesk CRM; blocks training conflicts; reroutes chat loads.", body_style), 
            Paragraph("Wait times cut from 185s to 15s. Prevented <b>$3,200/day</b> in SLA penalties.", body_style)
        ],
        [
            Paragraph("<b>VIP Priority Routing</b>", body_style), 
            Paragraph("Scans SIP trunks via ACM server; identifies voice quality failures; prompts Tandem Care auth checkpoint.", body_style), 
            Paragraph("Escalated VIP calls to L3 in 10s. Retained <b>$24,000 VIP account contract ARR</b>.", body_style)
        ],
        [
            Paragraph("<b>Retention Discount Policy</b>", body_style), 
            Paragraph("Queries SharePoint policy documents; applies SOP-Finance-204 coupons; inserts updated records.", body_style), 
            Paragraph("Authorized $1,260 retention credit, saving <b>$8,400 contract value</b> (6.6x ROI relative to CAC).", body_style)
        ]
    ]
    
    t2 = Table(scenarios_data, colWidths=[1.5*inch, 2.5*inch, 2.5*inch])
    t2.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f3f4f6')),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('BOTTOMPADDING', (0,0), (-1,0), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#d1d5db')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
    ]))
    story.append(t2)
    
    doc.build(story)
    print("Executive Brief PDF generated: Avaya_Infinity_MCP_Executive_Brief.pdf")

if __name__ == "__main__":
    generate_pptx()
    generate_pdf()
